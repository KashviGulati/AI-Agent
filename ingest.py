"""
ingest.py - Document Processing and Storage System
Handles file management, content extraction, and vector storage
"""

import os
import io
import hashlib
import json
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Union

# Document processing libraries
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import docx
from pptx import Presentation

# RAG and ML libraries
import chromadb
from sentence_transformers import SentenceTransformer
from dotenv import load_dotenv

# Text processing
import nltk
from nltk.tokenize import sent_tokenize
import re

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# ---------- CONFIG ----------
CHROMA_DIR = "./chroma_db"
UPLOAD_DIR = "./uploaded_files"
METADATA_FILE = "./file_metadata.json"
COLLECTION_NAME = "document_store"
EMBED_MODEL = "all-MiniLM-L6-v2"
CHUNK_SIZE = 800
CHUNK_OVERLAP = 200
# ----------------------------

# Ensure directories exist
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(CHROMA_DIR, exist_ok=True)

# Load environment variables
load_dotenv()

class FileManager:
    """Manages file storage, metadata, and tracking"""
    
    def __init__(self, upload_dir: str, metadata_file: str):
        self.upload_dir = Path(upload_dir)
        self.metadata_file = Path(metadata_file)
        self.metadata = self._load_metadata()
    
    def _load_metadata(self) -> Dict:
        """Load file metadata from JSON file"""
        if self.metadata_file.exists():
            with open(self.metadata_file, 'r') as f:
                return json.load(f)
        return {}
    
    def _save_metadata(self):
        """Save metadata to JSON file"""
        with open(self.metadata_file, 'w') as f:
            json.dump(self.metadata, f, indent=2)
    
    def _get_file_hash(self, file_path: Union[str, Path]) -> str:
        """Generate MD5 hash of file content"""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def add_file(self, file_content: bytes, filename: str) -> Dict:
        """Add a new file to storage"""
        file_path = self.upload_dir / filename
        
        # Save file
        with open(file_path, 'wb') as f:
            f.write(file_content)
        
        # Calculate hash and metadata
        file_hash = self._get_file_hash(file_path)
        file_info = {
            "filename": filename,
            "file_path": str(file_path),
            "file_hash": file_hash,
            "upload_time": datetime.now().isoformat(),
            "file_size": len(file_content),
            "processed": False
        }
        
        # Check if file already exists with same hash
        for existing_id, existing_info in self.metadata.items():
            if existing_info.get("file_hash") == file_hash:
                return {"status": "duplicate", "file_id": existing_id, "message": "File already exists"}
        
        # Generate unique file ID
        file_id = f"{filename}_{file_hash[:8]}"
        self.metadata[file_id] = file_info
        self._save_metadata()
        
        return {"status": "added", "file_id": file_id, "file_info": file_info}
    
    def remove_file(self, file_id: str) -> Dict:
        """Remove file from storage and metadata"""
        if file_id not in self.metadata:
            return {"status": "not_found", "message": "File not found"}
        
        file_info = self.metadata[file_id]
        file_path = Path(file_info["file_path"])
        
        # Remove physical file
        if file_path.exists():
            file_path.unlink()
        
        # Remove from metadata
        del self.metadata[file_id]
        self._save_metadata()
        
        return {"status": "removed", "file_id": file_id}
    
    def list_files(self) -> Dict:
        """List all managed files"""
        return {"files": self.metadata}
    
    def mark_processed(self, file_id: str):
        """Mark file as processed"""
        if file_id in self.metadata:
            self.metadata[file_id]["processed"] = True
            self._save_metadata()
    
    def get_file_info(self, file_id: str) -> Optional[Dict]:
        """Get file information by ID"""
        return self.metadata.get(file_id)

class DocumentProcessor:
    """Enhanced document processing with better extraction"""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF with OCR fallback"""
        text = ""
        
        # Try direct text extraction first
        try:
            doc = fitz.open(file_path)
            for page in doc:
                page_text = page.get_text()
                if page_text.strip():
                    text += page_text + "\n"
            doc.close()
        except Exception as e:
            print(f"Direct PDF extraction failed: {e}")
        
        # If no text or very little text, use OCR
        if len(text.strip()) < 100:
            print(f"Using OCR for {file_path}...")
            try:
                doc = fitz.open(file_path)
                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    pix = page.get_pixmap(dpi=300)
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    ocr_text = pytesseract.image_to_string(img, lang='eng')
                    text += ocr_text + "\n"
                doc.close()
            except Exception as e:
                print(f"OCR failed: {e}")
                return "Error: Could not extract text from PDF"
        
        return text.strip()
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX files"""
        try:
            doc = docx.Document(file_path)
            text = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text.strip())
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text.append(" | ".join(row_text))
            
            return "\n".join(text)
            
        except Exception as e:
            return f"Error extracting DOCX: {str(e)}"
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== Slide {slide_num} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide header
                    text.extend(slide_text)
                    text.append("")  # Add blank line between slides
            
            return "\n".join(text)
            
        except Exception as e:
            return f"Error extracting PPTX: {str(e)}"
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from TXT files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
            except Exception as e:
                return f"Error reading TXT file: {str(e)}"
        except Exception as e:
            return f"Error extracting TXT: {str(e)}"
    
    @classmethod
    def extract_text(cls, file_path: str) -> str:
        """Main method to extract text based on file extension"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        extractors = {
            '.pdf': cls.extract_text_from_pdf,
            '.docx': cls.extract_text_from_docx,
            '.doc': cls.extract_text_from_docx,  # Assuming .doc can be read as .docx
            '.pptx': cls.extract_text_from_pptx,
            '.ppt': cls.extract_text_from_pptx,  # Assuming .ppt can be read as .pptx
            '.txt': cls.extract_text_from_txt,
        }
        
        extractor = extractors.get(extension)
        if not extractor:
            return f"Unsupported file type: {extension}"
        
        return extractor(str(file_path))

class TextChunker:
    """Intelligent text chunking with overlap"""
    
    @staticmethod
    def clean_text(text: str) -> str:
        """Clean and normalize text"""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)]', ' ', text)
        return text.strip()
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
        """Split text into overlapping chunks"""
        text = TextChunker.clean_text(text)
        
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        sentences = sent_tokenize(text)
        
        current_chunk = ""
        current_length = 0
        
        for sentence in sentences:
            sentence_length = len(sentence)
            
            if current_length + sentence_length > chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                
                # Create overlap by keeping last part of current chunk
                overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
                current_chunk = overlap_text + " " + sentence
                current_length = len(current_chunk)
            else:
                current_chunk += " " + sentence
                current_length += sentence_length + 1
        
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks

class DocumentIngestion:
    """Main document ingestion system"""
    
    def __init__(self):
        self.file_manager = FileManager(UPLOAD_DIR, METADATA_FILE)
        self.processor = DocumentProcessor()
        self.chunker = TextChunker()
        
        # Initialize ChromaDB
        self.client = chromadb.PersistentClient(path=CHROMA_DIR)
        try:
            self.collection = self.client.get_collection(name=COLLECTION_NAME)
        except Exception:
            self.collection = self.client.create_collection(name=COLLECTION_NAME)
        
        # Initialize embedding model
        self.embedder = None
        self._load_embedder()
    
    def _load_embedder(self):
        """Load embedding model"""
        if self.embedder is None:
            print("Loading embedding model...")
            self.embedder = SentenceTransformer(EMBED_MODEL)
            print("Embedding model loaded successfully!")
    
    def ingest_document(self, file_content: bytes, filename: str) -> Dict:
        """Main method to ingest a document"""
        print(f"\n=== Starting ingestion for: {filename} ===")
        
        # Add file to file manager
        result = self.file_manager.add_file(file_content, filename)
        
        if result["status"] == "duplicate":
            print(f"File {filename} already exists in the system.")
            return result
        
        file_id = result["file_id"]
        file_info = result["file_info"]
        
        try:
            # Step 1: Extract text
            print("Step 1: Extracting text...")
            text = self.processor.extract_text(file_info["file_path"])
            
            if not text or text.startswith("Error"):
                print(f"Text extraction failed: {text}")
                self.file_manager.remove_file(file_id)
                return {"status": "error", "message": f"Failed to extract text: {text}"}
            
            print(f"Extracted {len(text)} characters of text")
            
            # Step 2: Chunk text
            print("Step 2: Chunking text...")
            chunks = self.chunker.chunk_text(text)
            
            if not chunks:
                print("No text chunks created")
                self.file_manager.remove_file(file_id)
                return {"status": "error", "message": "No text chunks created"}
            
            print(f"Created {len(chunks)} chunks")
            
            # Step 3: Generate embeddings
            print("Step 3: Generating embeddings...")
            embeddings = self.embedder.encode(chunks, show_progress_bar=True).tolist()
            
            # Step 4: Store in ChromaDB
            print("Step 4: Storing in vector database...")
            documents = chunks
            metadatas = []
            ids = []
            
            for i, chunk in enumerate(chunks):
                chunk_id = f"{file_id}_chunk_{i}"
                ids.append(chunk_id)
                metadatas.append({
                    "file_id": file_id,
                    "filename": filename,
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "upload_time": file_info["upload_time"],
                    "chunk_size": len(chunk)
                })
            
            # Add to ChromaDB
            self.collection.add(
                embeddings=embeddings,
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
            
            # Mark as processed
            self.file_manager.mark_processed(file_id)
            
            print(f"✅ Successfully ingested {filename}")
            print(f"   - File ID: {file_id}")
            print(f"   - Chunks: {len(chunks)}")
            print(f"   - Total characters: {len(text)}")
            
            return {
                "status": "success",
                "file_id": file_id,
                "filename": filename,
                "chunks_created": len(chunks),
                "total_characters": len(text),
                "message": f"Document processed successfully with {len(chunks)} chunks"
            }
            
        except Exception as e:
            print(f"❌ Ingestion failed: {str(e)}")
            # Remove file if processing failed
            self.file_manager.remove_file(file_id)
            return {"status": "error", "message": f"Processing failed: {str(e)}"}
    
    def remove_document(self, file_id: str) -> Dict:
        """Remove document from the system"""
        print(f"\n=== Removing document: {file_id} ===")
        
        # Remove from ChromaDB
        try:
            # Get all chunk IDs for this file
            results = self.collection.get(
                where={"file_id": file_id},
                include=["metadatas"]
            )
            
            if results and results["ids"]:
                chunk_ids = results["ids"]
                self.collection.delete(ids=chunk_ids)
                print(f"Removed {len(chunk_ids)} chunks from vector database")
            else:
                print("No chunks found in vector database")
        except Exception as e:
            print(f"Error removing from ChromaDB: {e}")
        
        # Remove from file manager
        result = self.file_manager.remove_file(file_id)
        
        if result["status"] == "removed":
            print(f"✅ Successfully removed document: {file_id}")
        else:
            print(f"❌ Failed to remove document: {result.get('message', 'Unknown error')}")
        
        return result
    
    def bulk_ingest_from_directory(self, directory_path: str) -> Dict:
        """Ingest all supported files from a directory"""
        directory_path = Path(directory_path)
        
        if not directory_path.exists():
            return {"status": "error", "message": "Directory does not exist"}
        
        supported_extensions = {'.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt'}
        files_to_process = []
        
        for file_path in directory_path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                files_to_process.append(file_path)
        
        if not files_to_process:
            return {"status": "warning", "message": "No supported files found in directory"}
        
        print(f"\n=== Bulk ingestion from: {directory_path} ===")
        print(f"Found {len(files_to_process)} files to process")
        
        results = {
            "processed": [],
            "failed": [],
            "duplicates": [],
            "total_files": len(files_to_process)
        }
        
        for file_path in files_to_process:
            try:
                with open(file_path, 'rb') as f:
                    file_content = f.read()
                
                result = self.ingest_document(file_content, file_path.name)
                
                if result["status"] == "success":
                    results["processed"].append({
                        "filename": file_path.name,
                        "file_id": result["file_id"],
                        "chunks": result["chunks_created"]
                    })
                elif result["status"] == "duplicate":
                    results["duplicates"].append(file_path.name)
                else:
                    results["failed"].append({
                        "filename": file_path.name,
                        "error": result["message"]
                    })
                    
            except Exception as e:
                results["failed"].append({
                    "filename": file_path.name,
                    "error": str(e)
                })
        
        print(f"\n=== Bulk ingestion complete ===")
        print(f"Processed: {len(results['processed'])}")
        print(f"Failed: {len(results['failed'])}")
        print(f"Duplicates: {len(results['duplicates'])}")
        
        return {"status": "complete", "results": results}
    
    def get_ingestion_stats(self) -> Dict:
        """Get ingestion system statistics"""
        file_count = len(self.file_manager.metadata)
        processed_count = sum(1 for f in self.file_manager.metadata.values() if f.get("processed", False))
        
        try:
            collection_count = self.collection.count()
        except:
            collection_count = 0
        
        # Calculate total file size
        total_size = sum(f.get("file_size", 0) for f in self.file_manager.metadata.values())
        
        return {
            "total_files": file_count,
            "processed_files": processed_count,
            "pending_files": file_count - processed_count,
            "total_chunks": collection_count,
            "total_file_size_bytes": total_size,
            "total_file_size_mb": round(total_size / (1024*1024), 2),
            "supported_formats": [".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt"],
            "chunk_size": CHUNK_SIZE,
            "chunk_overlap": CHUNK_OVERLAP
        }
    
    def list_files(self) -> Dict:
        """List all files with detailed information"""
        return self.file_manager.list_files()

# CLI interface for standalone usage
if __name__ == "__main__":
    import argparse
    
    ingestion = DocumentIngestion()
    
    parser = argparse.ArgumentParser(description="Document Ingestion System")
    parser.add_argument("--ingest", type=str, help="Ingest a single file")
    parser.add_argument("--bulk", type=str, help="Bulk ingest from directory")
    parser.add_argument("--remove", type=str, help="Remove file by ID")
    parser.add_argument("--list", action="store_true", help="List all files")
    parser.add_argument("--stats", action="store_true", help="Show system statistics")
    
    args = parser.parse_args()
    
    if args.ingest:
        file_path = Path(args.ingest)
        if file_path.exists():
            with open(file_path, 'rb') as f:
                content = f.read()
            result = ingestion.ingest_document(content, file_path.name)
            print(json.dumps(result, indent=2))
        else:
            print("File not found!")
    
    elif args.bulk:
        result = ingestion.bulk_ingest_from_directory(args.bulk)
        print(json.dumps(result, indent=2))
    
    elif args.remove:
        result = ingestion.remove_document(args.remove)
        print(json.dumps(result, indent=2))
    
    elif args.list:
        result = ingestion.list_files()
        print(json.dumps(result, indent=2))
    
    elif args.stats:
        result = ingestion.get_ingestion_stats()
        print(json.dumps(result, indent=2))
    
    else:
        parser.print_help()